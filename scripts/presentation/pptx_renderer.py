"""PPTX Renderer — generates PowerPoint presentations from PresentationModel.

Renders structured slide data into a Telenor-branded .pptx file using the
official template (reference/pdf/Telenor_mal.pptx). Falls back to a blank
16:9 presentation when the template is unavailable.
"""

from __future__ import annotations

import logging
import os
import re
from typing import List, Optional, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

from .telegrafen import COLORS, PPTX_LAYOUT_MAP, TEMPLATE_PPTX, PROJECT_ROOT
from .slide_parser import PresentationModel, SlideModel

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Telenor PPTX colors (RGBColor instances)
# ---------------------------------------------------------------------------

TG_DARK_BLUE = RGBColor(7, 4, 82)
TG_MID_BLUE = RGBColor(28, 22, 197)
TG_TELENOR_BLUE = RGBColor(0, 200, 255)
TG_BLUE = RGBColor(41, 84, 255)
TG_LIGHT_CYAN = RGBColor(235, 255, 255)
WHITE = RGBColor(255, 255, 255)

# ---------------------------------------------------------------------------
# Font constant
# ---------------------------------------------------------------------------

FONT_PPTX = "Telenor Evolution PPT"

# ---------------------------------------------------------------------------
# Slide dimension constants (16:9 widescreen)
# ---------------------------------------------------------------------------

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Standard positioning
TITLE_LEFT = Inches(1)
TITLE_TOP = Inches(0.5)
TITLE_WIDTH = Inches(11)
TITLE_HEIGHT = Inches(1.2)

CONTENT_LEFT = Inches(1)
CONTENT_TOP = Inches(2)
CONTENT_WIDTH = Inches(11.333)
CONTENT_HEIGHT = Inches(4.5)

TABLE_LEFT = Inches(1)
TABLE_TOP = Inches(2)
TABLE_WIDTH = Inches(11.333)

# ---------------------------------------------------------------------------
# Layout categories (for content population strategy)
# ---------------------------------------------------------------------------

STATEMENT_LAYOUTS = {"hero-statement", "big-number", "key-finding"}

TWO_COLUMN_LAYOUTS = {
    "two-column", "dual-analysis", "comparison",
    "before-after", "context-split",
}

THREE_COLUMN_LAYOUTS = {"three-column"}

TITLE_ONLY_LAYOUTS = {"title-only", "divider", "chart", "animated-reveal"}

# ---------------------------------------------------------------------------
# Markdown stripping patterns
# ---------------------------------------------------------------------------

_MD_BOLD_RE = re.compile(r"\*\*(.+?)\*\*")
_MD_ITALIC_RE = re.compile(r"\*(.+?)\*")
_MD_LINK_RE = re.compile(r"\[([^\]]+)\]\([^)]+\)")
_MD_IMAGE_RE = re.compile(r"!\[([^\]]*)\]\([^)]+\)")
_MD_HEADING_RE = re.compile(r"^#{1,6}\s+", re.MULTILINE)
_MD_BULLET_RE = re.compile(r"^[\s]*[-*+]\s+", re.MULTILINE)
_MD_NUMBERED_RE = re.compile(r"^[\s]*\d+\.\s+", re.MULTILINE)
_MD_BLOCKQUOTE_RE = re.compile(r"^>\s*", re.MULTILINE)
_MD_CODE_BLOCK_RE = re.compile(r"```[\s\S]*?```")
_MD_INLINE_CODE_RE = re.compile(r"`([^`]+)`")
_MD_HTML_COMMENT_RE = re.compile(r"<!--.*?-->", re.DOTALL)
_MD_TABLE_LINE_RE = re.compile(r"^\|.*\|$", re.MULTILINE)
_MD_TABLE_SEP_RE = re.compile(r"^\|[\s:|-]+\|$", re.MULTILINE)


# ---------------------------------------------------------------------------
# Helper functions
# ---------------------------------------------------------------------------


def _get_layout(prs: Presentation, name: str) -> object:
    """Find a slide layout by name in the template.

    Tries an exact match first. If not found, attempts a case-insensitive
    match. Falls back to the first layout with 'Title' in its name, and
    ultimately to the very first layout in the template.

    Args:
        prs: An open python-pptx Presentation.
        name: The layout name to search for.

    Returns:
        The matching SlideLayout object.
    """
    # Exact match
    for layout in prs.slide_layouts:
        if layout.name == name:
            return layout

    # Case-insensitive match
    name_lower = name.lower()
    for layout in prs.slide_layouts:
        if layout.name.lower() == name_lower:
            logger.debug(
                "Layout '%s' matched case-insensitively as '%s'",
                name, layout.name,
            )
            return layout

    # Fallback: "Title and Content"
    for layout in prs.slide_layouts:
        if layout.name == "Title and Content":
            logger.warning(
                "Layout '%s' not found; falling back to 'Title and Content'",
                name,
            )
            return layout

    # Last resort: first layout
    logger.warning(
        "Layout '%s' not found and no 'Title and Content' available; "
        "using first layout '%s'",
        name, prs.slide_layouts[0].name,
    )
    return prs.slide_layouts[0]


def _delete_all_slides(prs: Presentation) -> None:
    """Remove every slide from *prs* using XML manipulation.

    python-pptx does not expose a public delete-slide API, so we
    manipulate the slide-ID list and drop the corresponding
    relationships directly.
    """
    ns = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
    while len(prs.slides._sldIdLst) > 0:
        rId = prs.slides._sldIdLst[0].attrib[f"{ns}id"]
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]


def _set_font(
    paragraph,
    name: Optional[str] = None,
    size: Optional[int] = None,
    bold: Optional[bool] = None,
    color: Optional[RGBColor] = None,
    alignment: Optional[int] = None,
) -> None:
    """Configure font properties on a paragraph's run(s).

    If *paragraph* has no runs, the settings are applied to the
    paragraph-level font. Otherwise they are applied to every existing run.

    Args:
        paragraph: A python-pptx Paragraph object.
        name: Font family name (e.g. ``"Telenor Evolution PPT"``).
        size: Font size as a ``Pt()`` value.
        bold: Whether the text should be bold.
        color: An ``RGBColor`` instance.
        alignment: A ``PP_ALIGN`` enum value.
    """
    font = paragraph.font
    if name:
        font.name = name
    if size:
        font.size = size
    if bold is not None:
        font.bold = bold
    if color:
        font.color.rgb = color
    if alignment is not None:
        paragraph.alignment = alignment


def _md_to_plain_text(md_text: str) -> str:
    """Strip markdown formatting to produce plain text suitable for PPTX.

    Handles bold, italic, links, images, headings, bullets, blockquotes,
    code blocks, inline code, and HTML comments.  Markdown tables are
    removed entirely (they are handled separately by the table builder).

    Args:
        md_text: Raw markdown string.

    Returns:
        Plain text with markdown syntax removed.
    """
    if not md_text:
        return ""

    text = md_text

    # Remove HTML comments (including slide/chart metadata blocks)
    text = _MD_HTML_COMMENT_RE.sub("", text)

    # Remove code blocks
    text = _MD_CODE_BLOCK_RE.sub("", text)

    # Remove images (keep alt text)
    text = _MD_IMAGE_RE.sub(r"\1", text)

    # Convert links to just their text
    text = _MD_LINK_RE.sub(r"\1", text)

    # Bold and italic
    text = _MD_BOLD_RE.sub(r"\1", text)
    text = _MD_ITALIC_RE.sub(r"\1", text)

    # Inline code
    text = _MD_INLINE_CODE_RE.sub(r"\1", text)

    # Headings — keep the text, drop the hashes
    text = _MD_HEADING_RE.sub("", text)

    # Blockquotes
    text = _MD_BLOCKQUOTE_RE.sub("", text)

    # Bullet markers — convert to simple dash
    text = _MD_BULLET_RE.sub("- ", text)

    # Remove markdown table rows (handled separately)
    text = _MD_TABLE_SEP_RE.sub("", text)
    text = _MD_TABLE_LINE_RE.sub("", text)

    # Collapse multiple blank lines
    text = re.sub(r"\n{3,}", "\n\n", text)

    return text.strip()


def _extract_markdown_table(md_text: str) -> Optional[List[List[str]]]:
    """Parse the first markdown table found in *md_text*.

    Returns a list of rows, each row being a list of cell strings.
    The header separator row (``|---|---|``) is omitted.
    Returns ``None`` if no table is found.
    """
    if not md_text:
        return None

    lines = md_text.split("\n")
    table_lines: list[str] = []
    in_table = False

    for line in lines:
        stripped = line.strip()
        if re.match(r"^\|.+\|$", stripped):
            in_table = True
            table_lines.append(stripped)
        elif in_table:
            # End of table block
            break

    if len(table_lines) < 2:
        return None

    rows: list[list[str]] = []
    for line in table_lines:
        # Skip separator rows  |---|---|
        if re.match(r"^\|[\s:|-]+\|$", line):
            continue
        cells = [c.strip() for c in line.split("|")]
        # split on | gives empty strings at start/end
        cells = [c for c in cells if c or cells.index(c) not in (0, len(cells) - 1)]
        # Cleaner: strip leading/trailing empty entries
        if cells and cells[0] == "":
            cells = cells[1:]
        if cells and cells[-1] == "":
            cells = cells[:-1]
        if cells:
            rows.append(cells)

    return rows if rows else None


def _body_without_table(md_text: str) -> str:
    """Return the body markdown with table lines removed."""
    if not md_text:
        return ""

    lines = md_text.split("\n")
    result: list[str] = []
    in_table = False

    for line in lines:
        stripped = line.strip()
        if re.match(r"^\|.+\|$", stripped):
            in_table = True
            continue
        elif in_table and not stripped:
            in_table = False
            continue
        else:
            in_table = False
            result.append(line)

    return "\n".join(result).strip()


def _split_content_columns(text: str, num_columns: int) -> List[str]:
    """Split body text into *num_columns* parts for multi-column layouts.

    Uses ``---`` or ``***`` horizontal-rule separators if present,
    otherwise splits by double-newline sections as evenly as possible.

    Args:
        text: Plain text body content.
        num_columns: Target number of columns (2 or 3).

    Returns:
        A list of strings, one per column. May be shorter than
        *num_columns* if not enough content is found.
    """
    # Try explicit separators first
    parts = re.split(r"\n\s*(?:---|\*\*\*|___)\s*\n", text)
    if len(parts) >= num_columns:
        return [p.strip() for p in parts[:num_columns]]

    # Fall back to splitting by paragraphs
    paragraphs = [p.strip() for p in text.split("\n\n") if p.strip()]
    if len(paragraphs) <= num_columns:
        # Pad with empty strings if needed
        while len(paragraphs) < num_columns:
            paragraphs.append("")
        return paragraphs

    # Distribute paragraphs evenly
    chunk_size = len(paragraphs) // num_columns
    columns: list[str] = []
    for i in range(num_columns):
        start = i * chunk_size
        if i == num_columns - 1:
            end = len(paragraphs)
        else:
            end = start + chunk_size
        columns.append("\n\n".join(paragraphs[start:end]))

    return columns


def _find_placeholder(slide, idx: int):
    """Find a placeholder on the slide by its index.

    Returns ``None`` if the placeholder does not exist.
    """
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


def _find_content_placeholders(slide) -> list:
    """Return all content placeholders (idx >= 1), sorted by index."""
    placeholders = []
    for ph in slide.placeholders:
        if ph.placeholder_format.idx >= 1:
            placeholders.append(ph)
    placeholders.sort(key=lambda p: p.placeholder_format.idx)
    return placeholders


def _add_table_to_slide(
    slide,
    table_data: List[List[str]],
    left: int = TABLE_LEFT,
    top: int = TABLE_TOP,
    width: int = TABLE_WIDTH,
) -> None:
    """Create a styled PPTX table from parsed markdown table data.

    The first row is treated as the header and receives a dark blue
    background with white text. Subsequent rows alternate between
    white and light-cyan backgrounds.

    Args:
        slide: The python-pptx slide to add the table to.
        table_data: List of rows, each a list of cell strings.
        left: Left position (EMU).
        top: Top position (EMU).
        width: Table width (EMU).
    """
    if not table_data:
        return

    num_rows = len(table_data)
    num_cols = max(len(row) for row in table_data)

    # Calculate row height: distribute available space
    available_height = Inches(4.5)
    row_height = int(available_height / num_rows)

    graphic = slide.shapes.add_table(
        num_rows, num_cols, left, top, width, available_height,
    )
    table = graphic.table

    # Distribute column widths evenly
    col_width = int(width / num_cols)
    for col_idx in range(num_cols):
        table.columns[col_idx].width = col_width

    for row_idx, row_data in enumerate(table_data):
        for col_idx in range(num_cols):
            cell = table.cell(row_idx, col_idx)
            cell_text = row_data[col_idx] if col_idx < len(row_data) else ""

            # Strip any remaining markdown from cell text
            cell_text = _MD_BOLD_RE.sub(r"\1", cell_text)
            cell_text = _MD_ITALIC_RE.sub(r"\1", cell_text)
            cell_text = _MD_LINK_RE.sub(r"\1", cell_text)

            cell.text = cell_text
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            # Style the paragraph
            for paragraph in cell.text_frame.paragraphs:
                _set_font(
                    paragraph,
                    name=FONT_PPTX,
                    size=Pt(12),
                    bold=(row_idx == 0),
                    color=WHITE if row_idx == 0 else TG_DARK_BLUE,
                )

            # Cell background
            cell_fill = cell.fill
            cell_fill.solid()
            if row_idx == 0:
                cell_fill.fore_color.rgb = TG_DARK_BLUE
            elif row_idx % 2 == 1:
                cell_fill.fore_color.rgb = RGBColor(255, 255, 255)
            else:
                cell_fill.fore_color.rgb = TG_LIGHT_CYAN


def _populate_title(slide, slide_model: SlideModel) -> None:
    """Set the title text on a slide.

    Tries placeholder idx 0 first. If no title placeholder exists,
    creates a text box in the standard title position.
    """
    title_text = slide_model.headline or slide_model.title
    if not title_text:
        return

    title_ph = _find_placeholder(slide, 0)
    if title_ph is not None:
        title_ph.text = title_text
        for paragraph in title_ph.text_frame.paragraphs:
            _set_font(
                paragraph,
                name=FONT_PPTX,
                size=Pt(28),
                bold=True,
            )
    else:
        # No title placeholder — add a text box
        txBox = slide.shapes.add_textbox(
            TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        _set_font(
            p,
            name=FONT_PPTX,
            size=Pt(28),
            bold=True,
            color=WHITE,
            alignment=PP_ALIGN.LEFT,
        )


def _populate_statement(slide, slide_model: SlideModel) -> None:
    """Populate a statement/hero layout with large display text.

    For hero-statement, big-number, and key-finding layouts the title
    placeholder receives oversized text. The body text (if any) goes
    into a subtitle area.
    """
    title_text = slide_model.headline or slide_model.title
    if not title_text:
        return

    title_ph = _find_placeholder(slide, 0)
    if title_ph is not None:
        title_ph.text = title_text
        for paragraph in title_ph.text_frame.paragraphs:
            font_size = Pt(44) if slide_model.layout == "big-number" else Pt(36)
            _set_font(
                paragraph,
                name=FONT_PPTX,
                size=font_size,
                bold=True,
            )
    else:
        txBox = slide.shapes.add_textbox(
            TITLE_LEFT, Inches(1.5), TITLE_WIDTH, Inches(2),
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        font_size = Pt(44) if slide_model.layout == "big-number" else Pt(36)
        _set_font(
            p,
            name=FONT_PPTX,
            size=font_size,
            bold=True,
            color=WHITE,
            alignment=PP_ALIGN.LEFT,
        )

    # Body text as subtitle
    body_text = _md_to_plain_text(slide_model.body_md)
    if not body_text:
        return

    content_phs = _find_content_placeholders(slide)
    if content_phs:
        ph = content_phs[0]
        ph.text = body_text
        for paragraph in ph.text_frame.paragraphs:
            _set_font(
                paragraph,
                name=FONT_PPTX,
                size=Pt(18),
            )
    else:
        txBox = slide.shapes.add_textbox(
            TITLE_LEFT, Inches(4), TITLE_WIDTH, Inches(2),
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = body_text
        _set_font(
            p,
            name=FONT_PPTX,
            size=Pt(18),
            color=TG_TELENOR_BLUE,
            alignment=PP_ALIGN.LEFT,
        )


def _populate_two_column(slide, slide_model: SlideModel) -> None:
    """Populate a two-column layout by splitting body content."""
    _populate_title(slide, slide_model)

    body_text = _md_to_plain_text(slide_model.body_md)
    if not body_text:
        return

    columns = _split_content_columns(body_text, 2)
    content_phs = _find_content_placeholders(slide)

    for i, col_text in enumerate(columns):
        if i < len(content_phs):
            ph = content_phs[i]
            ph.text = col_text
            for paragraph in ph.text_frame.paragraphs:
                _set_font(paragraph, name=FONT_PPTX, size=Pt(14))
        else:
            # Create a text box for the column
            col_left = CONTENT_LEFT + Inches(i * 5.667)
            col_width = Inches(5.333)
            txBox = slide.shapes.add_textbox(
                col_left, CONTENT_TOP, col_width, CONTENT_HEIGHT,
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = col_text
            _set_font(p, name=FONT_PPTX, size=Pt(14), color=TG_DARK_BLUE)


def _populate_three_column(slide, slide_model: SlideModel) -> None:
    """Populate a three-column layout by splitting body content."""
    _populate_title(slide, slide_model)

    body_text = _md_to_plain_text(slide_model.body_md)
    if not body_text:
        return

    columns = _split_content_columns(body_text, 3)
    content_phs = _find_content_placeholders(slide)

    for i, col_text in enumerate(columns):
        if i < len(content_phs):
            ph = content_phs[i]
            ph.text = col_text
            for paragraph in ph.text_frame.paragraphs:
                _set_font(paragraph, name=FONT_PPTX, size=Pt(13))
        else:
            col_left = CONTENT_LEFT + Inches(i * 3.778)
            col_width = Inches(3.444)
            txBox = slide.shapes.add_textbox(
                col_left, CONTENT_TOP, col_width, CONTENT_HEIGHT,
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = col_text
            _set_font(p, name=FONT_PPTX, size=Pt(13), color=TG_DARK_BLUE)


def _populate_standard(slide, slide_model: SlideModel) -> None:
    """Populate a standard single-content layout (briefing, evidence, etc.).

    Handles both body text and embedded markdown tables.
    """
    _populate_title(slide, slide_model)

    # Check for tables first
    table_data = _extract_markdown_table(slide_model.body_md)
    if table_data:
        _add_table_to_slide(slide, table_data)
        # Also populate any remaining non-table body text
        remaining = _body_without_table(slide_model.body_md)
        remaining_text = _md_to_plain_text(remaining)
        if remaining_text:
            content_phs = _find_content_placeholders(slide)
            if content_phs:
                ph = content_phs[0]
                ph.text = remaining_text
                for paragraph in ph.text_frame.paragraphs:
                    _set_font(paragraph, name=FONT_PPTX, size=Pt(14))
        return

    # Plain body text
    body_text = _md_to_plain_text(slide_model.body_md)
    if not body_text:
        return

    content_phs = _find_content_placeholders(slide)
    if content_phs:
        ph = content_phs[0]
        _fill_text_frame(ph.text_frame, body_text)
    else:
        txBox = slide.shapes.add_textbox(
            CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, CONTENT_HEIGHT,
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        _fill_text_frame(tf, body_text)


def _fill_text_frame(tf, body_text: str) -> None:
    """Fill a text frame with body text, preserving paragraph breaks.

    Each non-empty line in *body_text* becomes a paragraph. Lines
    starting with ``-`` are formatted as bullet points.
    """
    tf.word_wrap = True
    lines = body_text.split("\n")
    first = True

    for line in lines:
        line_stripped = line.strip()
        if not line_stripped:
            continue

        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        # Detect bullet lines
        is_bullet = line_stripped.startswith("- ")
        if is_bullet:
            p.text = line_stripped[2:]
            p.level = 0
        else:
            p.text = line_stripped

        _set_font(p, name=FONT_PPTX, size=Pt(14), color=TG_DARK_BLUE)
        p.space_before = Pt(4)
        p.space_after = Pt(2)


def _populate_title_only(slide, slide_model: SlideModel) -> None:
    """Populate title-only and divider layouts."""
    title_text = slide_model.headline or slide_model.title
    if not title_text:
        return

    title_ph = _find_placeholder(slide, 0)
    if title_ph is not None:
        title_ph.text = title_text
        for paragraph in title_ph.text_frame.paragraphs:
            _set_font(
                paragraph,
                name=FONT_PPTX,
                size=Pt(36),
                bold=True,
            )
    else:
        txBox = slide.shapes.add_textbox(
            TITLE_LEFT, Inches(2.5), TITLE_WIDTH, Inches(2),
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        _set_font(
            p,
            name=FONT_PPTX,
            size=Pt(36),
            bold=True,
            color=WHITE,
            alignment=PP_ALIGN.CENTER,
        )


def _populate_cards(slide, slide_model: SlideModel) -> None:
    """Populate cards-grid and celebration layouts.

    Each card becomes a text box positioned horizontally across
    the slide. Limited to 4 cards before wrapping.
    """
    _populate_title(slide, slide_model)

    cards = slide_model.cards
    if not cards:
        # Fall back to standard content if no structured cards
        body_text = _md_to_plain_text(slide_model.body_md)
        if body_text:
            content_phs = _find_content_placeholders(slide)
            if content_phs:
                content_phs[0].text = body_text
                for paragraph in content_phs[0].text_frame.paragraphs:
                    _set_font(paragraph, name=FONT_PPTX, size=Pt(14))
        return

    num_cards = min(len(cards), 4)
    card_width = Inches(10) / num_cards
    card_gap = Inches(0.333)
    start_left = CONTENT_LEFT

    for i, card in enumerate(cards[:4]):
        left = start_left + int(i * (card_width + card_gap))
        txBox = slide.shapes.add_textbox(
            left, CONTENT_TOP, int(card_width), Inches(4),
        )
        tf = txBox.text_frame
        tf.word_wrap = True

        # Card title
        p = tf.paragraphs[0]
        p.text = card.title or f"Card {i + 1}"
        _set_font(
            p,
            name=FONT_PPTX,
            size=Pt(16),
            bold=True,
            color=TG_TELENOR_BLUE,
        )

        # Card body
        if card.body:
            p2 = tf.add_paragraph()
            p2.text = _md_to_plain_text(card.body)
            _set_font(p2, name=FONT_PPTX, size=Pt(12), color=WHITE)
            p2.space_before = Pt(8)

        # Icon label (PPTX cannot embed SVG directly)
        if card.icon:
            p_icon = tf.add_paragraph()
            p_icon.text = f"[{card.icon}]"
            _set_font(
                p_icon,
                name=FONT_PPTX,
                size=Pt(10),
                color=TG_TELENOR_BLUE,
            )
            p_icon.space_before = Pt(12)


def _populate_numbers_panel(slide, slide_model: SlideModel) -> None:
    """Populate dashboard / numbers-panel layouts.

    Extracts numbers and labels from body text, placing each
    in its own text box.
    """
    _populate_title(slide, slide_model)

    body_text = _md_to_plain_text(slide_model.body_md)
    if not body_text:
        return

    # Try to use template content placeholders
    content_phs = _find_content_placeholders(slide)
    if content_phs:
        content_phs[0].text = body_text
        for paragraph in content_phs[0].text_frame.paragraphs:
            _set_font(paragraph, name=FONT_PPTX, size=Pt(14))
        return

    # Fallback: add as text box
    txBox = slide.shapes.add_textbox(
        CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, CONTENT_HEIGHT,
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    _fill_text_frame(tf, body_text)


def _add_speaker_notes(slide_obj, slide_model: SlideModel) -> None:
    """Attach speaker notes to the slide if present."""
    notes_text = slide_model.speaker_notes
    if not notes_text:
        return

    try:
        notes_slide = slide_obj.notes_slide
        notes_slide.notes_text_frame.text = _md_to_plain_text(notes_text)
    except Exception as exc:
        logger.warning(
            "Could not add speaker notes to slide '%s': %s",
            slide_model.title, exc,
        )


def _populate_content(slide, slide_model: SlideModel, prs: Presentation) -> None:
    """Fill slide placeholders based on layout type.

    Dispatches to specialized population functions depending on
    the slide's layout category.

    Args:
        slide: The python-pptx Slide object.
        slide_model: The parsed SlideModel with content.
        prs: The parent Presentation (for layout inspection).
    """
    layout = slide_model.layout

    if layout in STATEMENT_LAYOUTS:
        _populate_statement(slide, slide_model)
    elif layout in TWO_COLUMN_LAYOUTS:
        _populate_two_column(slide, slide_model)
    elif layout in THREE_COLUMN_LAYOUTS:
        _populate_three_column(slide, slide_model)
    elif layout in TITLE_ONLY_LAYOUTS:
        _populate_title_only(slide, slide_model)
    elif layout in ("cards-grid", "celebration"):
        _populate_cards(slide, slide_model)
    elif layout in ("dashboard", "numbers-panel"):
        _populate_numbers_panel(slide, slide_model)
    elif layout == "icon-showcase":
        # Icon-showcase gets title + card-like content
        _populate_cards(slide, slide_model)
    elif layout == "table":
        _populate_standard(slide, slide_model)
    else:
        # Default: standard single-content layout
        _populate_standard(slide, slide_model)


# ---------------------------------------------------------------------------
# Main renderer
# ---------------------------------------------------------------------------


def render_pptx(presentation: PresentationModel, output_path: str) -> None:
    """Generate a PowerPoint file from a PresentationModel.

    Opens the Telenor PPTX template, clears any existing slides, and
    builds new slides from the ``presentation.slides`` list. Each slide
    is matched to a named layout in the template, populated with its
    title, body text, tables, and speaker notes, then saved to
    *output_path*.

    If the template file does not exist, a blank 16:9 presentation
    is created as a fallback.

    Args:
        presentation: The parsed PresentationModel to render.
        output_path: Filesystem path for the generated ``.pptx`` file.

    Raises:
        OSError: If *output_path* cannot be written to.
    """
    # Open template or create blank fallback
    if os.path.exists(TEMPLATE_PPTX):
        prs = Presentation(TEMPLATE_PPTX)
        _delete_all_slides(prs)
        logger.info("Loaded template: %s", TEMPLATE_PPTX)
    else:
        logger.warning(
            "Template not found at %s — creating blank 16:9 presentation",
            TEMPLATE_PPTX,
        )
        prs = Presentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT

    # Log available layouts for debugging
    layout_names = [layout.name for layout in prs.slide_layouts]
    logger.debug("Available layouts (%d): %s", len(layout_names), layout_names)

    # Build each slide
    for slide_model in presentation.slides:
        # Skip notes-only slides (paired as speaker notes on parent)
        if slide_model.slide_id.lower().endswith("b") and slide_model.speaker_notes is None:
            # This is a speaker-notes slide that was already paired
            # Check if it was consumed by _pair_speaker_notes
            parent_id = slide_model.slide_id.lower()[:-1]
            parent_exists = any(
                s.slide_id.lower() == parent_id for s in presentation.slides
            )
            if parent_exists:
                logger.debug(
                    "Skipping notes slide '%s' (paired with '%s')",
                    slide_model.slide_id, parent_id,
                )
                continue

        # Determine layout name
        if slide_model.pptx_layout:
            layout_name = slide_model.pptx_layout
        else:
            layout_name = PPTX_LAYOUT_MAP.get(
                slide_model.layout, "Title and Content",
            )

        layout = _get_layout(prs, layout_name)
        logger.debug(
            "Slide %d (%s): layout='%s' -> '%s'",
            slide_model.number,
            slide_model.slide_id,
            slide_model.layout,
            layout.name,
        )

        # Add the slide
        slide_obj = prs.slides.add_slide(layout)

        # Populate content
        _populate_content(slide_obj, slide_model, prs)

        # Add speaker notes
        _add_speaker_notes(slide_obj, slide_model)

    # Ensure output directory exists
    output_dir = os.path.dirname(output_path)
    if output_dir:
        os.makedirs(output_dir, exist_ok=True)

    # Save
    prs.save(output_path)
    logger.info(
        "Saved PPTX with %d slides to %s",
        len(prs.slides), output_path,
    )
